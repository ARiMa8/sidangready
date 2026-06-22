from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Iterable

SCAN_LIKE_PDF_WARNING = (
    "Dokumen berhasil diunggah, tetapi teks yang dapat dibaca sangat sedikit. "
    "Kemungkinan file berupa scan/gambar. OCR belum tersedia pada versi MVP."
)
LOW_TEXT_WARNING = (
    "Dokumen berhasil dibaca, tetapi teks yang dapat diekstrak sangat sedikit. "
    "Periksa kembali apakah file berisi teks yang dapat dipilih."
)


class DocumentParserError(Exception):
    """Raised when a document cannot be parsed for text extraction."""


@dataclass(frozen=True)
class ParsedDocument:
    text: str
    extraction_status: str
    extraction_warning: str | None = None
    page_count: int | None = None
    slide_count: int | None = None


def parse_document_bytes(
    *,
    file_name: str,
    file_mime_type: str,
    file_bytes: bytes,
    document_type: str,
) -> ParsedDocument:
    extension = Path(file_name).suffix.lower()

    if extension == ".pdf" or file_mime_type == "application/pdf":
        return _parse_pdf(file_bytes, document_type)
    if (
        extension == ".docx"
        or file_mime_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        return _parse_docx(file_bytes)
    if (
        extension == ".pptx"
        or file_mime_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        return _parse_pptx(file_bytes)
    if extension == ".txt" or file_mime_type == "text/plain":
        return _parse_txt(file_bytes)

    raise DocumentParserError("Format dokumen belum didukung untuk ekstraksi teks.")


def _normalize_text(parts: Iterable[str]) -> str:
    normalized_parts = []
    for part in parts:
        cleaned = "\n".join(line.rstrip() for line in part.replace("\r\n", "\n").split("\n"))
        cleaned = cleaned.strip()
        if cleaned:
            normalized_parts.append(cleaned)
    return "\n\n".join(normalized_parts).strip()


def _quality_status(
    text: str,
    *,
    warning: str = LOW_TEXT_WARNING,
    page_count: int | None = None,
) -> tuple[str, str | None]:
    readable_chars = len(text.strip())
    if readable_chars < 300:
        return "low_quality", warning
    if page_count and page_count > 0 and readable_chars / page_count < 120:
        return "low_quality", warning
    return "success", None


def _parse_pdf(file_bytes: bytes, document_type: str) -> ParsedDocument:
    try:
        import fitz
    except ImportError as exc:
        raise DocumentParserError(
            "PyMuPDF belum terpasang. Jalankan instalasi dependency backend."
        ) from exc

    try:
        document = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        raise DocumentParserError("PDF tidak dapat dibuka untuk ekstraksi teks.") from exc

    pages: list[str] = []
    for page_index, page in enumerate(document, start=1):
        page_text = page.get_text("text").strip()
        if page_text:
            pages.append(f"[Halaman {page_index}]\n{page_text}")

    text = _normalize_text(pages)
    page_count = document.page_count
    status, warning = _quality_status(
        text,
        warning=SCAN_LIKE_PDF_WARNING,
        page_count=page_count,
    )
    slide_count = page_count if document_type == "slides" else None

    return ParsedDocument(
        text=text,
        extraction_status=status,
        extraction_warning=warning,
        page_count=page_count,
        slide_count=slide_count,
    )


def _parse_docx(file_bytes: bytes) -> ParsedDocument:
    try:
        from docx import Document as WordDocument
    except ImportError as exc:
        raise DocumentParserError(
            "python-docx belum terpasang. Jalankan instalasi dependency backend."
        ) from exc

    try:
        document = WordDocument(BytesIO(file_bytes))
    except Exception as exc:
        raise DocumentParserError("DOCX tidak dapat dibuka untuk ekstraksi teks.") from exc

    parts: list[str] = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    text = _normalize_text(parts)
    status, warning = _quality_status(text)
    return ParsedDocument(
        text=text,
        extraction_status=status,
        extraction_warning=warning,
    )


def _parse_pptx(file_bytes: bytes) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentParserError(
            "python-pptx belum terpasang. Jalankan instalasi dependency backend."
        ) from exc

    try:
        presentation = Presentation(BytesIO(file_bytes))
    except Exception as exc:
        raise DocumentParserError("PPTX tidak dapat dibuka untuk ekstraksi teks.") from exc

    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = list(_extract_slide_text_parts(slide.shapes))
        slide_text = _normalize_text(slide_text_parts)
        title = slide_text.splitlines()[0] if slide_text else f"Slide {slide_index}"
        if slide_text:
            slides.append(f"[Slide {slide_index}: {title}]\n{slide_text}")
        else:
            slides.append(f"[Slide {slide_index}: Tanpa teks terbaca]")

    text = _normalize_text(slides)
    status, warning = _quality_status(text)
    return ParsedDocument(
        text=text,
        extraction_status=status,
        extraction_warning=warning,
        slide_count=len(presentation.slides),
    )


def _extract_slide_text_parts(shapes) -> Iterable[str]:
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                yield text

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    yield row_text

        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _extract_slide_text_parts(nested_shapes)


def _parse_txt(file_bytes: bytes) -> ParsedDocument:
    last_error: UnicodeDecodeError | None = None
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            text = file_bytes.decode(encoding)
            normalized_text = _normalize_text([text])
            status, warning = _quality_status(normalized_text)
            return ParsedDocument(
                text=normalized_text,
                extraction_status=status,
                extraction_warning=warning,
            )
        except UnicodeDecodeError as exc:
            last_error = exc

    raise DocumentParserError("TXT tidak dapat dibaca sebagai teks.") from last_error
