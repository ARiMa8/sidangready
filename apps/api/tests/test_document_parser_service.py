from __future__ import annotations

from io import BytesIO
import importlib.util
import unittest

from app.services.document_parser_service import (
    SCAN_LIKE_PDF_WARNING,
    parse_document_bytes,
)


class DocumentParserServiceTestCase(unittest.TestCase):
    def test_parses_txt_bytes(self) -> None:
        parsed = parse_document_bytes(
            file_name="catatan-revisi.txt",
            file_mime_type="text/plain",
            file_bytes=(
                "Catatan revisi: perjelas batasan penelitian dan sesuaikan "
                "kesimpulan dengan tujuan penelitian. "
            ).encode("utf-8")
            * 5,
            document_type="revision_notes",
        )

        self.assertEqual(parsed.extraction_status, "success")
        self.assertIn("Catatan revisi", parsed.text)

    @unittest.skipUnless(importlib.util.find_spec("docx"), "python-docx is not installed")
    def test_parses_docx_bytes(self) -> None:
        from docx import Document

        document = Document()
        document.add_paragraph("Judul Skripsi")
        document.add_paragraph(
            "Isi laporan skripsi yang cukup panjang untuk melewati ambang "
            "kualitas ekstraksi dokumen. "
            * 5
        )
        buffer = BytesIO()
        document.save(buffer)

        parsed = parse_document_bytes(
            file_name="laporan.docx",
            file_mime_type=(
                "application/vnd.openxmlformats-officedocument."
                "wordprocessingml.document"
            ),
            file_bytes=buffer.getvalue(),
            document_type="thesis",
        )

        self.assertEqual(parsed.extraction_status, "success")
        self.assertIn("Judul Skripsi", parsed.text)

    @unittest.skipUnless(importlib.util.find_spec("pptx"), "python-pptx is not installed")
    def test_parses_pptx_bytes(self) -> None:
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "Metodologi Pengujian"
        slide.placeholders[1].text = (
            "Slide menjelaskan pengujian black-box dan evaluasi model. " * 6
        )
        buffer = BytesIO()
        presentation.save(buffer)

        parsed = parse_document_bytes(
            file_name="slides.pptx",
            file_mime_type=(
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            ),
            file_bytes=buffer.getvalue(),
            document_type="slides",
        )

        self.assertEqual(parsed.extraction_status, "success")
        self.assertEqual(parsed.slide_count, 1)
        self.assertIn("[Slide 1: Metodologi Pengujian]", parsed.text)

    @unittest.skipUnless(importlib.util.find_spec("fitz"), "PyMuPDF is not installed")
    def test_parses_pdf_bytes_and_warns_for_low_quality_pdf(self) -> None:
        import fitz

        document = fitz.open()
        page = document.new_page()
        page.insert_text((72, 72), "Sedikit")
        pdf_bytes = document.tobytes()

        parsed = parse_document_bytes(
            file_name="scan-like.pdf",
            file_mime_type="application/pdf",
            file_bytes=pdf_bytes,
            document_type="thesis",
        )

        self.assertEqual(parsed.extraction_status, "low_quality")
        self.assertEqual(parsed.extraction_warning, SCAN_LIKE_PDF_WARNING)
        self.assertEqual(parsed.page_count, 1)
